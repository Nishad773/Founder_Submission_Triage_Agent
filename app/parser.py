from __future__ import annotations

from io import BytesIO

import fitz
from pptx import Presentation

from app.schemas import DocumentBundle
from app.utils import DOCUMENT_LABELS, ensure_supported_file, prepare_text_for_llm


def _split_pdf_page(page: fitz.Page) -> list[str]:
    blocks = page.get_text("blocks")
    lines: list[str] = []
    for block in blocks:
        for line in block[4].splitlines():
            if line.strip():
                lines.append(line)
    return lines


def extract_pdf_text(file_bytes: bytes) -> str:
    lines: list[str] = []
    with fitz.open(stream=file_bytes, filetype="pdf") as document:
        for page_number, page in enumerate(document, start=1):
            page_lines = _split_pdf_page(page)
            if page_lines:
                lines.append(f"Page {page_number}")
                lines.extend(page_lines)
    return prepare_text_for_llm(lines)


def extract_pptx_text(file_bytes: bytes) -> str:
    presentation = Presentation(BytesIO(file_bytes))
    lines: list[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_lines.extend(shape.text.splitlines())
        if slide_lines:
            lines.append(f"Slide {slide_number}")
            lines.extend(slide_lines)
    return prepare_text_for_llm(lines)


def extract_text(file_bytes: bytes, extension: str) -> str:
    if extension == ".pdf":
        return extract_pdf_text(file_bytes)
    if extension == ".pptx":
        return extract_pptx_text(file_bytes)
    raise ValueError("Unsupported file type. Only PDF and PPTX are allowed.")


def extract_file_text(filename: str, file_bytes: bytes) -> str:
    return extract_text(file_bytes, ensure_supported_file(filename))


def build_document_bundle(files_by_category: dict[str, list[tuple[str, bytes]]]) -> DocumentBundle:
    bundle_data: dict[str, str] = {key: "" for key in DOCUMENT_LABELS}

    for category, uploads in files_by_category.items():
        if category not in bundle_data or not uploads:
            continue

        sections: list[str] = []
        for filename, file_bytes in uploads:
            extracted_text = extract_file_text(filename, file_bytes)
            if extracted_text:
                source_name = DOCUMENT_LABELS[category]
                sections.append(f"[{source_name} | {filename}]\n{extracted_text}")

        bundle_data[category] = "\n\n".join(sections)

    return DocumentBundle.model_validate(bundle_data)
